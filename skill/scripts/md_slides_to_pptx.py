#!/usr/bin/env python3
"""Convert Marp-ish Markdown slides (--- separated) to PPTX with speaker notes.

Usage:
  python md_slides_to_pptx.py --input slides.md --out deck.pptx

Speaker notes: HTML comments containing a line `SPEAKER: ...` (see paper-ppt.md).

Requires: pip install python-pptx
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

_SPEAKER_RE = re.compile(
    r"<!--([\s\S]*?)-->",
    re.MULTILINE,
)
_SPEAKER_LINE = re.compile(r"^\s*SPEAKER:\s*(.*)\s*$", re.IGNORECASE | re.MULTILINE)


def _strip_front_matter(text: str) -> str:
    if text.startswith("---"):
        end = text.find("\n---", 3)
        if end != -1:
            return text[end + 4 :].lstrip("\n")
    return text


def _split_slides(text: str) -> list[str]:
    body = _strip_front_matter(text)
    parts = re.split(r"\n---\s*\n", body)
    return [p.strip() for p in parts if p.strip()]


def _extract_notes(block: str) -> tuple[str, str]:
    notes: list[str] = []

    def _collect(m: re.Match[str]) -> str:
        inner = m.group(1)
        for sm in _SPEAKER_LINE.finditer(inner):
            notes.append(sm.group(1).strip())
        # keep SOURCE lines out of visible slide body
        return ""

    cleaned = _SPEAKER_RE.sub(_collect, block)
    return cleaned.strip(), "\n".join(notes).strip()


def _slide_title_and_bullets(block: str) -> tuple[str, list[str]]:
    lines = [ln.rstrip() for ln in block.splitlines() if ln.strip()]
    if not lines:
        return "Slide", []
    title = lines[0].lstrip("# ").strip()
    bullets: list[str] = []
    for ln in lines[1:]:
        s = ln.strip()
        if s.startswith(("-", "*", "•")):
            bullets.append(s.lstrip("-*• ").strip())
        elif s.startswith("#"):
            break
        else:
            bullets.append(s)
    return title or "Slide", bullets[:8]


def build_pptx(slides: list[str], out: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as e:
        raise SystemExit(
            "python-pptx is required. Install with: pip install python-pptx"
        ) from e

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    notes_n = 0

    for block in slides:
        visible, notes = _extract_notes(block)
        title, bullets = _slide_title_and_bullets(visible)
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.0))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True

        body = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.8), Inches(5.5))
        btf = body.text_frame
        btf.word_wrap = True
        if not bullets:
            btf.paragraphs[0].text = ""
        for i, item in enumerate(bullets):
            para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            para.text = item
            para.level = 0
            para.font.size = Pt(18)
            para.space_after = Pt(8)

        if notes:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = notes
            notes_n += 1

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote {out} ({len(slides)} slides, {notes_n} with notes)")


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--input", "-i", required=True, help="Marp-ish slides.md")
    ap.add_argument("--out", "-o", required=True, help="Output .pptx path")
    args = ap.parse_args(argv)

    src = Path(args.input)
    if not src.is_file():
        print(f"Input not found: {src}", file=sys.stderr)
        return 1
    slides = _split_slides(src.read_text(encoding="utf-8"))
    if not slides:
        print("No slides found (use --- separators).", file=sys.stderr)
        return 1
    build_pptx(slides, Path(args.out))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
